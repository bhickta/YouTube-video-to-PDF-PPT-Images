from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches 
import os


def sortfiles(imagefiles):
    sorted_list1 = []
    for image in imagefiles:
        sorted_list1.append(int(image.split('.')[0]))
    sorted_list1.sort()
    
    sorted_list2 = []
    for element in sorted_list1:
        sorted_list2.append(str(element)+'.jpeg')
    return sorted_list2


def images2pdf(name, location):
    WIDTH = 210
    HEIGHT = 297    
    pdf = FPDF(orientation = 'Landscape', unit = 'mm', format='A4')

    for image in sortfiles(os.listdir(location)):
        pdf.add_page()
        pdf.image(location+str(image), x=10, y=10, w = HEIGHT-10)

    pdf.output(str(name)+'.pdf')
    print("Conversion successful! " + str(name)+'.pdf'+ " ---> Stored at "+ location)


def images2ppt(name, location):
    """https://www.geeksforgeeks.org/creating-and-updating-powerpoint-presentations-in-python-using-python-pptx/"""
    # Creating presentation object  
    ppt = Presentation()
    # Creating slide layout
    blank_slide_layout= ppt.slide_layouts[6]

    for image in sortfiles(os.listdir(location)):
        # Creating slide object to add in ppt i.e. adding slides to ppt
        slide = ppt.slides.add_slide(blank_slide_layout) 
        # add a picture shape to slide
        slide.shapes.add_picture(location+str(image), left=0, top=0, width=ppt.slide_width, height=ppt.slide_height)

    ppt.save(str(name)+'.pptx')
    print("Conversion successful! " + str(name)+'.pptx'+ " ---> Stored at "+ location)
